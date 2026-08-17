#!/usr/bin/env python3
"""Create the native Google Slides deck for the superset-security-check demo video.

The deck content (on-screen bullets, speaker-notes narration and timing hints) lives
in ``DECK`` below, so the companion narration script in
``docs/superset-security-check-script.md`` is generated from the same source of truth.

Usage:

    # Regenerate only the markdown narration script (no credentials needed)
    python scripts/generate_deck.py --script-only

    # Build a .pptx (no credentials needed) to import into Google Slides manually
    python scripts/generate_deck.py --pptx build/superset-security-check.pptx

    # Create the native Google Slides presentation and print its shareable link
    python scripts/generate_deck.py --share-with you@example.com

Credentials (one of):

    GOOGLE_APPLICATION_CREDENTIALS   path to a service-account JSON
    GOOGLE_SERVICE_ACCOUNT_JSON      the service-account JSON itself
    --credentials PATH               path to a service-account JSON
    GOOGLE_OAUTH_CLIENT_SECRETS      OAuth client secrets JSON (+ GOOGLE_OAUTH_TOKEN
                                     for a cached authorized-user token)

Required scopes: https://www.googleapis.com/auth/presentations and
https://www.googleapis.com/auth/drive.
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

SCOPES = [
    "https://www.googleapis.com/auth/presentations",
    "https://www.googleapis.com/auth/drive",
]

PLAYBOOK_ID = "7adc898234024db386a1c66082489651"
PLAYBOOK_URL = f"https://app.devin.ai/org/kumar-gaurav-demo/settings/playbooks/{PLAYBOOK_ID}"

PRESENTER = "Kumar Gaurav"
PRESENTATION_DATE = "17th Aug 2026"
DECK_TITLE = "superset-security-check — Devin security review on every push"

REPO_ROOT = Path(__file__).resolve().parent.parent
SCRIPT_DOC = REPO_ROOT / "docs" / "superset-security-check-script.md"

# Colors: dark text on light backgrounds only.
INK = {"red": 0.07, "green": 0.09, "blue": 0.15}
MUTED = {"red": 0.35, "green": 0.40, "blue": 0.47}
BOX_FILL = {"red": 0.93, "green": 0.95, "blue": 0.98}
BOX_LINE = {"red": 0.20, "green": 0.29, "blue": 0.37}
ACCENT_FILL = {"red": 0.99, "green": 0.95, "blue": 0.86}
ACCENT_LINE = {"red": 0.72, "green": 0.45, "blue": 0.09}

# --------------------------------------------------------------------------------------
# Deck content
# --------------------------------------------------------------------------------------

DIAGRAM_NODES: List[Tuple[str, float, float]] = [
    # (label, x, y) — 200x55pt boxes on a 720x405pt slide, read as a serpentine.
    ("1. Git push", 40, 105),
    ("2. POST /webhook/commit", 270, 105),
    ("3. HMAC-SHA256 verify\n(401 if invalid)", 500, 105),
    ("4. extract_commits()\nnormalize payload", 500, 190),
    ("5. build_devin_session_payload()\nprompt + playbook_id", 270, 190),
    ("6. POST /organizations/{org}/sessions\nDevin Sessions API", 40, 190),
    ("7. SQLite ledger\nrecord_session()", 40, 275),
    ("8. /metrics (JSON)\n+ /dashboard (HTML)", 270, 275),
]

PLAYBOOK_NODE = (
    f"Playbook {PLAYBOOK_ID}\n(DEVIN_PLAYBOOK_ID)",
    500,
    275,
)

# (from_x, from_y, to_x, to_y, dashed, label) — the arrowhead sits on the "to" end.
DIAGRAM_ARROWS: List[Tuple[float, float, float, float, bool, str]] = [
    (240, 132, 270, 132, False, ""),
    (470, 132, 500, 132, False, ""),
    (600, 160, 600, 190, False, ""),
    (500, 217, 470, 217, False, ""),
    (270, 217, 240, 217, False, ""),
    (140, 245, 140, 275, False, ""),
    (240, 302, 270, 302, False, ""),
    (600, 275, 380, 248, False, ""),  # playbook -> build_devin_session_payload
    (300, 277, 140, 250, True, "poll-status/{id}"),
]

DECK: List[Dict[str, Any]] = [
    {
        "kind": "title",
        "title": "superset-security-check",
        "subtitle": (
            "Turning every Git push into an autonomous Devin security review\n"
            f"{PRESENTER} — {PRESENTATION_DATE}"
        ),
        "timing": "0:00–0:10",
        "beat": "Title",
        "notes": (
            "Hi, I'm Kumar Gaurav. This is superset-security-check: a small service that "
            "turns every Git push into an autonomous Devin security review, and shows the "
            "result on a live dashboard."
        ),
    },
    {
        "kind": "bullets",
        "title": "What — the problem",
        "bullets": [
            "Every push is a risk window: leaked secrets, insecure dependencies, subtle logic flaws",
            "Security review lands too late — manual PR review, or after ship",
            "Teams want review on every push, in every repo, but human bandwidth doesn't scale",
            "This project hands every push to Devin: automatic, immediate, traceable",
        ],
        "timing": "0:10–0:55",
        "beat": "What",
        "notes": (
            "Every push is a risk window. A commit can leak a credential, pull in a vulnerable "
            "dependency, or quietly introduce a logic flaw no linter will notice. In most teams "
            "the security review happens far too late — in a manual pull request review, if "
            "someone has time, or after the code has shipped. What teams want is a review on "
            "every push, in every repo, every time, and that is exactly where human bandwidth "
            "runs out. So this project puts the one thing that does scale — an autonomous agent "
            "— directly on the push event. superset-security-check receives a push webhook, "
            "verifies it, and hands the change to Devin: automatic, immediate, and recorded."
        ),
        "source": "docs/README.md",
    },
    {
        "kind": "diagram",
        "title": "How — end-to-end architecture",
        "timing": "0:55–1:25",
        "beat": "How",
        "callout": (
            f"Playbook input — DEVIN_PLAYBOOK_ID = {PLAYBOOK_ID}  ·  {PLAYBOOK_URL}  ·  "
            "effect: added to the session request body as playbook_id"
        ),
        "notes": (
            "Here is the whole system, and the thing to notice is how little of it is mine. A "
            "push hits the webhook. The signature is verified. The payload is normalized into a "
            "commit summary. That summary becomes a prompt, combined with a playbook ID, and "
            "posted to the Devin Sessions API. The session is written to a SQLite ledger, which "
            "is aggregated into metrics at /metrics and a live dashboard. That's it — thin, "
            "boring plumbing. Devin is the intelligence; this is just the wiring that gets a "
            "real-world event to it."
        ),
        "source": "docs/architecture.md",
    },
    {
        "kind": "bullets",
        "title": "How — ingestion & validation",
        "bullets": [
            "Webhook endpoint: POST /webhook/commit",
            "HMAC-SHA256 signature verified with a constant-time compare",
            "Forged or unsigned requests get 401 — nothing downstream runs",
        ],
        "timing": "1:25–1:40",
        "beat": "How",
        "notes": (
            "Ingestion first. Pushes arrive at /webhook/commit. Every request's HMAC-SHA256 "
            "signature is recomputed from the shared secret and compared in constant time, so a "
            "forged request gets a 401 and never reaches Devin. Trust boundary first, work "
            "second."
        ),
        "source": "app.py verify_signature, webhook.py commit_webhook",
    },
    {
        "kind": "bullets",
        "title": "How — payload normalization",
        "bullets": [
            "Raw push flattened to: repository, branch, head commit, commit count, messages",
            "One small summary object instead of a provider-shaped payload",
            "Decouples the service from GitHub / GitLab / Bitbucket payload differences",
        ],
        "timing": "1:40–1:55",
        "beat": "How",
        "notes": (
            "Then normalization. The raw push is flattened into a small summary: repository, "
            "branch, head commit, commit count, messages. That one step decouples everything "
            "downstream from the provider's payload shape — GitHub, GitLab and Bitbucket all "
            "look the same from here."
        ),
        "source": "app.py extract_commits",
    },
    {
        "kind": "bullets",
        "title": "How — Devin orchestration  +  DEMO",
        "bullets": [
            "Summary → natural-language prompt → POST /organizations/{org}/sessions",
            "Optional DEVIN_PLAYBOOK_ID injects the standardized review procedure",
            "DEMO: push a commit, watch the session appear live on /dashboard",
        ],
        "code": (
            "{\n"
            '  "name": "Webhook commit - repo:branch",\n'
            '  "prompt": "Repository: code-lgtm/superset\n'
            "              Branch: main\n"
            "              Recent commit messages:\n"
            '              - fix(daos): reject ...",\n'
            f'  "playbook_id":\n      "{PLAYBOOK_ID}"\n'
            "}"
        ),
        "timing": "1:55–2:25",
        "beat": "How",
        "notes": (
            "Now orchestration. The summary becomes a natural-language prompt, posted to the "
            "Devin Sessions API — one session per push. If DEVIN_PLAYBOOK_ID is set it rides "
            "along as playbook_id, so every session follows the same vetted review procedure. "
            "Let me show it: I push a commit, GitHub fires the webhook, the signature checks "
            "out, and a Devin session spins up and starts reviewing that diff. Refresh the "
            "dashboard and there it is. No human touched anything in between."
        ),
        "source": "webhook.py build_devin_session_payload, create_devin_session",
    },
    {
        "kind": "bullets",
        "title": "How — persistence & monitoring",
        "bullets": [
            "Every session logged to the SQLite ledger",
            "Aggregated metrics: success rate, throughput, per-repo breakdown, duration percentiles",
            "Served as JSON at /metrics and as an auto-refreshing /dashboard",
            "Honest labelling: result = session created, not a task verdict",
            "Status refreshed on demand via /poll-status/<session_id>",
        ],
        "timing": "2:25–2:45",
        "beat": "How",
        "notes": (
            "Every session is written to the SQLite ledger and aggregated into metrics — success "
            "rate, throughput, a per-repository breakdown and duration percentiles — exposed as "
            "JSON and as an auto-refreshing dashboard. One honest caveat: result records that "
            "the session was created, not a verdict on the task; live status comes from the "
            "poll-status endpoint."
        ),
        "source": "analytics.py record_session, get_session_metrics, poll_devin_session_status",
    },
    {
        "kind": "bullets",
        "title": "Why Devin is uniquely suited",
        "bullets": [
            "A static CI scanner applies fixed rules — it only finds what it was pre-programmed to find",
            "Devin reasons openly per commit: reads the diff, traces data flow, checks dependencies",
            "Scales to every push in every repo — one session per push, zero humans, no fatigue",
            "Playbooks standardize it: the same vetted procedure, every single time",
            "The plumbing is simple; the intelligence is Devin",
        ],
        "timing": "2:45–4:00",
        "beat": "Why",
        "notes": (
            "So why an agent, and not the scanner you already have in CI? A static scanner "
            "applies fixed rules. It is fast and useful, but it only ever finds what someone "
            "pre-programmed it to look for; it has no idea what this commit was trying to do. "
            "Devin does something categorically different: it reasons openly about the change in "
            "front of it. It reads the actual diff, traces the data flow through the functions "
            "that changed, asks whether user input reaches something dangerous, checks the new "
            "dependencies, and explains what it found in plain language. That is judgement you "
            "would normally book a human reviewer for — and unlike a human it scales without "
            "limit: one session per push, every repo, at three in the morning, on the hundredth "
            "commit of the day, with the same attention as the first. Playbooks remove the "
            "remaining variance: the same playbook ID on every session means the same vetted "
            "procedure every time. And that is the real point here. The plumbing you just saw is "
            "deliberately "
            "simple — a few hundred lines of Flask, an HMAC check, a SQLite table. It is only "
            "worth building because the hard, open-ended part, actually reviewing the code, is "
            "done autonomously by the agent."
        ),
        "source": "webhook.py",
    },
    {
        "kind": "bullets",
        "title": "When — next steps in a real engagement",
        "bullets": [
            "Durable persistence: managed DB or persistent volume instead of ephemeral SQLite",
            "Close the feedback loop: post findings back to the PR as a review or status check",
            "Richer triggers & routing: PR events, per-repo prompts, Slack and ticketing",
            "Automated status reconciliation with a background poller",
            "Deeper analytics: trends over time, per-repo risk scoring",
        ],
        "timing": "4:00–4:40",
        "beat": "When",
        "notes": (
            "In a real engagement, here is the order I would do it in. First, durable "
            "persistence: the SQLite ledger is ephemeral on Cloud Run, so it moves to a managed "
            "database or a persistent volume. Second, close the feedback loop — post findings "
            "back onto the pull request as a review or status check, where the decision actually "
            "gets made. Third, richer triggers and routing: pull requests as well as pushes, "
            "per-repo prompts, Slack and ticketing. Fourth, a background poller to reconcile "
            "status automatically. Then deeper analytics — trends and per-repo risk scoring."
        ),
        "source": "README.md (Session ledger persistence), analytics.py",
    },
    {
        "kind": "closing",
        "title": "Closing",
        "body": (
            "A small, trustworthy service captures a real-world event and delegates the "
            "open-ended thinking to an autonomous agent.\n\n"
            "The plumbing is simple on purpose — the leverage comes from Devin."
        ),
        "timing": "4:40–4:45",
        "beat": "Closing",
        "notes": (
            "So: a small, trustworthy service captures a real-world event and delegates the "
            "open-ended thinking to an autonomous agent. The plumbing is simple on purpose — "
            "the leverage comes from Devin. Thank you."
        ),
    },
]

BEAT_ORDER = ["Title", "What", "How", "Why", "When", "Closing"]
BEAT_BUDGET = {
    "Title": "0:00–0:10 (~10s)",
    "What": "0:10–0:55 (~45s)",
    "How": "0:55–2:45 (~1m50s, includes the live demo)",
    "Why": "2:45–4:00 (~1m15s)",
    "When": "4:00–4:40 (~40s)",
    "Closing": "4:40–4:45 (~5s)",
}


# --------------------------------------------------------------------------------------
# Narration script (markdown companion)
# --------------------------------------------------------------------------------------


def _word_count(text: str) -> int:
    return len(text.split())


def build_script_markdown() -> str:
    total_words = sum(_word_count(slide["notes"]) for slide in DECK)
    lines: List[str] = [
        "# superset-security-check — video narration script",
        "",
        f"Presenter: **{PRESENTER}** · Date: **{PRESENTATION_DATE}** · Target runtime: "
        "**~4m45s** (under 5 minutes)",
        "",
        "Read the narration verbatim; the same text is in each slide's speaker notes. "
        "Generated by `scripts/generate_deck.py` — edit the `DECK` definition there, not this "
        "file.",
        "",
        f"Total narration: ~{total_words} words (~{total_words // 160}m"
        f"{round(total_words % 160 / 160 * 60):02d}s at a normal 160 words-per-minute pace).",
        "",
    ]

    for beat in BEAT_ORDER:
        slides = [s for s in DECK if s["beat"] == beat]
        if not slides:
            continue
        lines.append(f"## {beat} — {BEAT_BUDGET[beat]}")
        lines.append("")
        for slide in slides:
            index = DECK.index(slide) + 1
            lines.append(f"### Slide {index} — {slide['title']}  ·  {slide['timing']}")
            lines.append("")
            lines.append(slide["notes"])
            lines.append("")
            if slide.get("source"):
                lines.append(f"_Source: {slide['source']}_")
                lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def write_script_doc() -> Path:
    SCRIPT_DOC.parent.mkdir(parents=True, exist_ok=True)
    SCRIPT_DOC.write_text(build_script_markdown(), encoding="utf-8")
    return SCRIPT_DOC


# --------------------------------------------------------------------------------------
# Google API helpers
# --------------------------------------------------------------------------------------


class MissingCredentials(RuntimeError):
    pass


def load_credentials(credentials_path: Optional[str]):
    """Return Google credentials from a service account or an OAuth token."""
    from google.oauth2 import service_account
    from google.oauth2.credentials import Credentials as UserCredentials

    path = credentials_path or os.environ.get("GOOGLE_APPLICATION_CREDENTIALS", "")
    raw = os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON", "")
    info: Optional[Dict[str, Any]] = None

    if path and Path(path).is_file():
        info = json.loads(Path(path).read_text(encoding="utf-8"))
    elif raw.strip().startswith("{"):
        info = json.loads(raw)

    if info is not None:
        if info.get("type") == "service_account":
            return service_account.Credentials.from_service_account_info(info, scopes=SCOPES)
        if "refresh_token" in info:
            return UserCredentials.from_authorized_user_info(info, scopes=SCOPES)
        if "installed" in info or "web" in info:
            return _oauth_flow(info)

    token_raw = os.environ.get("GOOGLE_OAUTH_TOKEN", "")
    if token_raw.strip().startswith("{"):
        return UserCredentials.from_authorized_user_info(json.loads(token_raw), scopes=SCOPES)

    client_secrets = os.environ.get("GOOGLE_OAUTH_CLIENT_SECRETS", "")
    if client_secrets.strip().startswith("{"):
        return _oauth_flow(json.loads(client_secrets))
    if client_secrets and Path(client_secrets).is_file():
        return _oauth_flow(json.loads(Path(client_secrets).read_text(encoding="utf-8")))

    raise MissingCredentials(
        "No Google credentials found. Set GOOGLE_APPLICATION_CREDENTIALS (or "
        "GOOGLE_SERVICE_ACCOUNT_JSON) to a service-account JSON with the Slides and Drive "
        "scopes enabled, or provide GOOGLE_OAUTH_CLIENT_SECRETS / GOOGLE_OAUTH_TOKEN."
    )


def _oauth_flow(client_config: Dict[str, Any]):
    from google_auth_oauthlib.flow import InstalledAppFlow

    flow = InstalledAppFlow.from_client_config(client_config, scopes=SCOPES)
    return flow.run_local_server(port=0)


# --------------------------------------------------------------------------------------
# Slide building
# --------------------------------------------------------------------------------------


def _pt(value: float) -> Dict[str, Any]:
    return {"magnitude": value, "unit": "PT"}


def _element_properties(slide_id: str, x: float, y: float, w: float, h: float) -> Dict[str, Any]:
    return {
        "pageObjectId": slide_id,
        "size": {"width": _pt(w), "height": _pt(h)},
        "transform": {"scaleX": 1, "scaleY": 1, "translateX": x, "translateY": y, "unit": "PT"},
    }


def _text_style_request(object_id: str, size: float, bold: bool, color: Dict[str, float],
                        font: str = "Arial") -> Dict[str, Any]:
    return {
        "updateTextStyle": {
            "objectId": object_id,
            "style": {
                "fontFamily": font,
                "fontSize": _pt(size),
                "bold": bold,
                "foregroundColor": {"opaqueColor": {"rgbColor": color}},
            },
            "textRange": {"type": "ALL"},
            "fields": "fontFamily,fontSize,bold,foregroundColor",
        }
    }


def _box_requests(slide_id: str, object_id: str, text: str, x: float, y: float, w: float,
                  h: float, *, font_size: float = 9, fill: Dict[str, float] = BOX_FILL,
                  line: Dict[str, float] = BOX_LINE, bold: bool = True,
                  align: str = "CENTER") -> List[Dict[str, Any]]:
    return [
        {
            "createShape": {
                "objectId": object_id,
                "shapeType": "ROUND_RECTANGLE",
                "elementProperties": _element_properties(slide_id, x, y, w, h),
            }
        },
        {"insertText": {"objectId": object_id, "text": text}},
        {
            "updateShapeProperties": {
                "objectId": object_id,
                "shapeProperties": {
                    "shapeBackgroundFill": {"solidFill": {"color": {"rgbColor": fill}}},
                    "outline": {
                        "outlineFill": {"solidFill": {"color": {"rgbColor": line}}},
                        "weight": _pt(1),
                    },
                    "contentAlignment": "MIDDLE",
                },
                "fields": "shapeBackgroundFill.solidFill.color,outline.outlineFill.solidFill.color,outline.weight,contentAlignment",
            }
        },
        {
            "updateParagraphStyle": {
                "objectId": object_id,
                "style": {"alignment": align},
                "textRange": {"type": "ALL"},
                "fields": "alignment",
            }
        },
        _text_style_request(object_id, font_size, bold, INK),
    ]


def _arrow_requests(slide_id: str, object_id: str, x1: float, y1: float, x2: float, y2: float,
                    dashed: bool) -> List[Dict[str, Any]]:
    """Create an arrow between two points.

    A Slides line runs from the top-left to the bottom-right of its bounding box, so an arrow
    pointing up and/or left is drawn as the same segment with the arrowhead on its start.
    """
    dx, dy = x2 - x1, y2 - y1
    forward = (dx >= 0) and (dy >= 0)
    requests: List[Dict[str, Any]] = [
        {
            "createLine": {
                "objectId": object_id,
                "lineCategory": "STRAIGHT",
                "elementProperties": _element_properties(
                    slide_id, min(x1, x2), min(y1, y2), max(abs(dx), 1), max(abs(dy), 1)
                ),
            }
        }
    ]
    properties: Dict[str, Any] = {
        "lineFill": {"solidFill": {"color": {"rgbColor": BOX_LINE}}},
        "weight": _pt(1.5),
        "endArrow": "FILL_ARROW" if forward else "NONE",
        "startArrow": "NONE" if forward else "FILL_ARROW",
    }
    fields = "lineFill.solidFill.color,weight,endArrow,startArrow"
    if dashed:
        properties["dashStyle"] = "DASH"
        fields += ",dashStyle"
    requests.append({
        "updateLineProperties": {"objectId": object_id, "lineProperties": properties,
                                 "fields": fields}
    })
    return requests


def _label_requests(slide_id: str, object_id: str, text: str, x: float, y: float, w: float,
                    h: float, size: float = 8, color: Dict[str, float] = MUTED,
                    align: str = "START") -> List[Dict[str, Any]]:
    return [
        {
            "createShape": {
                "objectId": object_id,
                "shapeType": "TEXT_BOX",
                "elementProperties": _element_properties(slide_id, x, y, w, h),
            }
        },
        {"insertText": {"objectId": object_id, "text": text}},
        {
            "updateParagraphStyle": {
                "objectId": object_id,
                "style": {"alignment": align},
                "textRange": {"type": "ALL"},
                "fields": "alignment",
            }
        },
        _text_style_request(object_id, size, False, color),
    ]


def build_slide_requests() -> List[Dict[str, Any]]:
    """Requests that create every slide and its placeholder text."""
    requests: List[Dict[str, Any]] = []

    for index, slide in enumerate(DECK):
        slide_id = f"slide_{index + 1}"
        title_id = f"{slide_id}_title"
        body_id = f"{slide_id}_body"

        if slide["kind"] == "title":
            requests.append({
                "createSlide": {
                    "objectId": slide_id,
                    "insertionIndex": index,
                    "slideLayoutReference": {"predefinedLayout": "TITLE"},
                    "placeholderIdMappings": [
                        {"layoutPlaceholder": {"type": "CENTERED_TITLE", "index": 0},
                         "objectId": title_id},
                        {"layoutPlaceholder": {"type": "SUBTITLE", "index": 0},
                         "objectId": body_id},
                    ],
                }
            })
            requests.append({"insertText": {"objectId": title_id, "text": slide["title"]}})
            requests.append({"insertText": {"objectId": body_id, "text": slide["subtitle"]}})
            requests.append(_text_style_request(title_id, 40, True, INK))
            requests.append(_text_style_request(body_id, 16, False, INK))
        elif slide["kind"] == "diagram":
            requests.append({
                "createSlide": {
                    "objectId": slide_id,
                    "insertionIndex": index,
                    "slideLayoutReference": {"predefinedLayout": "TITLE_ONLY"},
                    "placeholderIdMappings": [
                        {"layoutPlaceholder": {"type": "TITLE", "index": 0},
                         "objectId": title_id},
                    ],
                }
            })
            requests.append({"insertText": {"objectId": title_id, "text": slide["title"]}})
            requests.append(_text_style_request(title_id, 24, True, INK))
            requests.extend(_diagram_requests(slide_id, slide["callout"]))
        elif slide["kind"] == "closing":
            requests.append({
                "createSlide": {
                    "objectId": slide_id,
                    "insertionIndex": index,
                    "slideLayoutReference": {"predefinedLayout": "TITLE_AND_BODY"},
                    "placeholderIdMappings": [
                        {"layoutPlaceholder": {"type": "TITLE", "index": 0},
                         "objectId": title_id},
                        {"layoutPlaceholder": {"type": "BODY", "index": 0},
                         "objectId": body_id},
                    ],
                }
            })
            requests.append({"insertText": {"objectId": title_id, "text": slide["title"]}})
            requests.append({"insertText": {"objectId": body_id, "text": slide["body"]}})
            requests.append(_text_style_request(title_id, 28, True, INK))
            requests.append(_text_style_request(body_id, 18, False, INK))
        else:
            requests.append({
                "createSlide": {
                    "objectId": slide_id,
                    "insertionIndex": index,
                    "slideLayoutReference": {"predefinedLayout": "TITLE_AND_BODY"},
                    "placeholderIdMappings": [
                        {"layoutPlaceholder": {"type": "TITLE", "index": 0},
                         "objectId": title_id},
                        {"layoutPlaceholder": {"type": "BODY", "index": 0},
                         "objectId": body_id},
                    ],
                }
            })
            requests.append({"insertText": {"objectId": title_id, "text": slide["title"]}})
            requests.append({
                "insertText": {"objectId": body_id, "text": "\n".join(slide["bullets"])}
            })
            requests.append({
                "createParagraphBullets": {
                    "objectId": body_id,
                    "textRange": {"type": "ALL"},
                    "bulletPreset": "BULLET_DISC_CIRCLE_SQUARE",
                }
            })
            requests.append(_text_style_request(title_id, 24, True, INK))
            requests.append(_text_style_request(body_id, 15, False, INK))
            if slide.get("code"):
                code_id = f"{slide_id}_code"
                requests.extend(_box_requests(
                    slide_id, code_id, slide["code"], 360, 250, 330, 110,
                    font_size=8, bold=False, align="START",
                ))
                requests.append({
                    "updateTextStyle": {
                        "objectId": code_id,
                        "style": {"fontFamily": "Courier New"},
                        "textRange": {"type": "ALL"},
                        "fields": "fontFamily",
                    }
                })

        footer_id = f"{slide_id}_footer"
        requests.extend(_label_requests(
            slide_id, footer_id, f"{slide['beat']}  ·  {slide['timing']}", 500, 375, 190, 20,
            size=8, align="END",
        ))

    return requests


def _diagram_requests(slide_id: str, callout: str) -> List[Dict[str, Any]]:
    requests: List[Dict[str, Any]] = []
    for i, (label, x, y) in enumerate(DIAGRAM_NODES):
        requests.extend(_box_requests(slide_id, f"{slide_id}_n{i}", label, x, y, 200, 55))

    label, x, y = PLAYBOOK_NODE
    requests.extend(_box_requests(
        slide_id, f"{slide_id}_playbook", label, x, y, 200, 55,
        fill=ACCENT_FILL, line=ACCENT_LINE, font_size=8,
    ))

    for i, (x1, y1, x2, y2, dashed, text) in enumerate(DIAGRAM_ARROWS):
        requests.extend(_arrow_requests(slide_id, f"{slide_id}_a{i}", x1, y1, x2, y2, dashed))
        if text:
            requests.extend(_label_requests(
                slide_id, f"{slide_id}_a{i}_label", text, (x1 + x2) / 2 - 30,
                (y1 + y2) / 2 - 14, 150, 14, size=7,
            ))

    requests.extend(_label_requests(
        slide_id, f"{slide_id}_callout", callout, 40, 345, 650, 40, size=8, color=INK,
    ))
    return requests


def notes_requests(slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Insert the narration into each slide's speaker-notes shape."""
    requests: List[Dict[str, Any]] = []
    for index, (api_slide, content) in enumerate(zip(slides, DECK)):
        notes_page = api_slide.get("slideProperties", {}).get("notesPage", {})
        notes_id = notes_page.get("notesProperties", {}).get("speakerNotesObjectId")
        if not notes_id:
            print(f"warning: slide {index + 1} has no speaker-notes shape", file=sys.stderr)
            continue
        text = (
            f"[{content['beat']} · {content['timing']}]\n\n{content['notes']}"
        )
        if content.get("source"):
            text += f"\n\nSource: {content['source']}"
        requests.append({"insertText": {"objectId": notes_id, "text": text}})
    return requests


# --------------------------------------------------------------------------------------
# Entry point
# --------------------------------------------------------------------------------------


# --------------------------------------------------------------------------------------
# PowerPoint rendering (import into Google Slides manually)
# --------------------------------------------------------------------------------------


def _rgb(color: Dict[str, float]):
    from pptx.dml.color import RGBColor

    return RGBColor(*(int(round(color[channel] * 255)) for channel in ("red", "green", "blue")))


def _pptx_textbox(slide, text: str, x: float, y: float, w: float, h: float, *, size: float,
                  bold: bool = False, color: Dict[str, float] = INK, align: str = "LEFT",
                  font: str = "Arial", space_after: float = 0):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    frame = box.text_frame
    frame.word_wrap = True
    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = getattr(PP_ALIGN, align)
        paragraph.space_after = Pt(space_after)
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.name = font
            run.font.color.rgb = _rgb(color)
    return box


def _pptx_box(slide, text: str, x: float, y: float, w: float, h: float, *, size: float = 9,
              fill: Dict[str, float] = BOX_FILL, line: Dict[str, float] = BOX_LINE,
              align: str = "CENTER", anchor: str = "MIDDLE", font: str = "Arial",
              bold_first: bool = True):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = getattr(MSO_ANCHOR, anchor)
    for index, part in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = part
        paragraph.alignment = getattr(PP_ALIGN, align)
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold_first and index == 0
            run.font.name = font
            run.font.color.rgb = _rgb(INK)
    return shape


def _pptx_arrow(slide, x1: float, y1: float, x2: float, y2: float, dashed: bool) -> None:
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Pt

    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
    connector.line.color.rgb = _rgb(BOX_LINE)
    connector.line.width = Pt(1.5)
    line_element = connector.line._get_or_add_ln()
    if dashed:
        dash = line_element.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash", {"val": "dash"}
        )
        line_element.append(dash)
    head = line_element.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd",
        {"type": "triangle", "w": "med", "len": "med"},
    )
    line_element.append(head)


def build_pptx(path: Path) -> Path:
    """Render the same deck as a .pptx for manual import into Google Slides."""
    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    presentation.slide_width = Pt(720)
    presentation.slide_height = Pt(405)
    blank_layout = presentation.slide_layouts[6]

    for slide_content in DECK:
        slide = presentation.slides.add_slide(blank_layout)
        kind = slide_content["kind"]

        if kind == "title":
            _pptx_textbox(slide, slide_content["title"], 60, 130, 600, 60, size=40, bold=True,
                          align="CENTER")
            _pptx_textbox(slide, slide_content["subtitle"], 60, 200, 600, 80, size=16,
                          align="CENTER")
        else:
            _pptx_textbox(slide, slide_content["title"], 40, 35, 650, 45, size=24, bold=True)

        if kind == "bullets":
            body = "\n".join(f"•  {bullet}" for bullet in slide_content["bullets"])
            width = 290 if slide_content.get("code") else 650
            _pptx_textbox(slide, body, 40, 100, width, 240, size=15, space_after=10)
            if slide_content.get("code"):
                _pptx_box(slide, slide_content["code"], 350, 100, 340, 130, size=9,
                          align="LEFT", anchor="TOP", font="Courier New", bold_first=False)
        elif kind == "closing":
            _pptx_textbox(slide, slide_content["body"], 40, 120, 650, 200, size=18)
        elif kind == "diagram":
            for label, x, y in DIAGRAM_NODES:
                _pptx_box(slide, label, x, y, 200, 55)
            label, x, y = PLAYBOOK_NODE
            _pptx_box(slide, label, x, y, 200, 55, size=8, fill=ACCENT_FILL, line=ACCENT_LINE)
            for x1, y1, x2, y2, dashed, text in DIAGRAM_ARROWS:
                _pptx_arrow(slide, x1, y1, x2, y2, dashed)  # arrowhead lands on (x2, y2)
                if text:
                    _pptx_textbox(slide, text, (x1 + x2) / 2 - 30, (y1 + y2) / 2 - 14, 150, 14,
                                  size=7, color=MUTED)
            _pptx_textbox(slide, slide_content["callout"], 40, 345, 650, 45, size=8)

        _pptx_textbox(slide, f"{slide_content['beat']}  ·  {slide_content['timing']}", 500, 378,
                      190, 18, size=8, color=MUTED, align="RIGHT")

        notes = f"[{slide_content['beat']} · {slide_content['timing']}]\n\n{slide_content['notes']}"
        if slide_content.get("source"):
            notes += f"\n\nSource: {slide_content['source']}"
        slide.notes_slide.notes_text_frame.text = notes

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))
    return path


def create_presentation(credentials, share_with: Optional[str], anyone_with_link: bool) -> str:
    from googleapiclient.discovery import build

    slides_api = build("slides", "v1", credentials=credentials)
    drive_api = build("drive", "v3", credentials=credentials)

    presentation = slides_api.presentations().create(body={"title": DECK_TITLE}).execute()
    presentation_id = presentation["presentationId"]
    default_slide_id = presentation["slides"][0]["objectId"]

    requests = build_slide_requests()
    requests.append({"deleteObject": {"objectId": default_slide_id}})
    slides_api.presentations().batchUpdate(
        presentationId=presentation_id, body={"requests": requests}
    ).execute()

    created = slides_api.presentations().get(presentationId=presentation_id).execute()
    notes = notes_requests(created["slides"])
    if notes:
        slides_api.presentations().batchUpdate(
            presentationId=presentation_id, body={"requests": notes}
        ).execute()

    if anyone_with_link:
        drive_api.permissions().create(
            fileId=presentation_id,
            body={"role": "writer", "type": "anyone"},
            supportsAllDrives=True,
        ).execute()
    if share_with:
        drive_api.permissions().create(
            fileId=presentation_id,
            body={"role": "writer", "type": "user", "emailAddress": share_with},
            sendNotificationEmail=False,
            supportsAllDrives=True,
        ).execute()

    return f"https://docs.google.com/presentation/d/{presentation_id}/edit"


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--credentials", help="path to a service-account JSON")
    parser.add_argument("--share-with", help="email address to grant edit access to")
    parser.add_argument(
        "--no-link-sharing",
        action="store_true",
        help="do not grant anyone-with-the-link access",
    )
    parser.add_argument(
        "--script-only",
        action="store_true",
        help="only regenerate docs/superset-security-check-script.md",
    )
    parser.add_argument(
        "--pptx",
        metavar="PATH",
        help="write the deck as a .pptx (for manual import into Google Slides) and exit",
    )
    args = parser.parse_args()

    path = write_script_doc()
    print(f"Wrote narration script: {path}")
    if args.script_only:
        return 0

    if args.pptx:
        print(f"Wrote PowerPoint deck: {build_pptx(Path(args.pptx))}")
        return 0

    try:
        credentials = load_credentials(args.credentials)
    except MissingCredentials as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2

    url = create_presentation(credentials, args.share_with, not args.no_link_sharing)
    print(f"Google Slides deck: {url}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
